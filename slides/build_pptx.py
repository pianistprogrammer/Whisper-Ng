#!/usr/bin/env python3
"""Assemble PNG slides into a 16:9 PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

for i in range(1, 14):
    png = f"/Users/I558118/Documents/Projects/Whisper-Ng/slides/slide{i:02d}.png"
    slide = prs.slides.add_slide(BLANK)
    slide.shapes.add_picture(png, 0, 0, prs.slide_width, prs.slide_height)
    print(f"Added slide {i:02d}")

out = "/Users/I558118/Documents/Projects/Whisper-Ng/Whisper_NG_Conference_Presentation.pptx"
prs.save(out)
print(f"\nSaved: {out}")
