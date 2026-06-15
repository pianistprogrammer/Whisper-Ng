#!/usr/bin/env python3
"""Generate conference presentation for Dr. Abimbola Jeremiah."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
GREEN_DARK  = RGBColor(0x00, 0x6B, 0x3C)   # Nigeria deep green
GREEN_MID   = RGBColor(0x00, 0x96, 0x55)   # medium green accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF4, 0xF4, 0xF4)
CHARCOAL    = RGBColor(0x22, 0x22, 0x22)
GOLD        = RGBColor(0xE8, 0xB8, 0x00)
LIGHT_GREEN = RGBColor(0xD6, 0xF0, 0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # totally blank


# ── Helper utilities ─────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=Pt(16), bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def section_header(slide, title):
    """Green left bar + slide title."""
    add_rect(slide, 0, 0, 0.12, 7.5, fill=GREEN_DARK)
    add_text(slide, title, 0.3, 0.18, 12.8, 0.7,
             font_size=Pt(30), bold=True, color=GREEN_DARK)
    add_rect(slide, 0.3, 0.95, 12.5, 0.04, fill=GREEN_MID)


def footer(slide, slide_num, total=12):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=GREEN_DARK)
    add_text(slide, "Dr. Abimbola Jeremiah  |  Silesian University of Technology",
             0.15, 7.12, 9, 0.3, font_size=Pt(10), color=WHITE)
    add_text(slide, f"{slide_num} / {total}", 12.5, 7.12, 0.7, 0.3,
             font_size=Pt(10), color=WHITE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=GREEN_DARK)            # full bg
add_rect(slide, 0, 0, 0.35, 7.5, fill=GOLD)                   # gold stripe

# White content card
add_rect(slide, 0.65, 0.8, 12.2, 5.6, fill=WHITE)

add_text(slide,
         'If It Doesn\'t Understand Me,\nIs It Really Intelligent?',
         0.85, 1.0, 11.8, 2.4,
         font_size=Pt(36), bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

add_rect(slide, 2.5, 3.55, 8.3, 0.05, fill=GOLD)

add_text(slide,
         'Fine-Tuning OpenAI Whisper for Nigerian Languages\n'
         'Hausa · Yoruba · Igbo · Nigerian Pidgin',
         0.85, 3.7, 11.8, 1.0,
         font_size=Pt(18), bold=False, color=CHARCOAL, align=PP_ALIGN.CENTER)

add_text(slide,
         'Dr. Abimbola Jeremiah\nAssistant Professor, Silesian University of Technology, Poland (Online)',
         0.85, 4.85, 11.8, 1.0,
         font_size=Pt(15), bold=True, color=GREEN_MID, align=PP_ALIGN.CENTER)

add_text(slide, 'June 2026', 0.85, 6.1, 11.8, 0.4,
         font_size=Pt(13), color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Presentation Overview")
footer(slide, 2)

items = [
    ("1", "The Spark — Why an App Became a Research Project"),
    ("2", "The Problem — Speech Recognition That Doesn't Know You"),
    ("3", "Technical Approach — Fine-Tuning Whisper with LoRA"),
    ("4", "Datasets — Google WaxalNLP, Mozilla Common Voice & More"),
    ("5", "Training Results & What the Numbers Tell Us"),
    ("6", "Challenges — Tones, Diacritics & Data Scarcity"),
    ("7", "Crowdsourcing Solution — ngs.primeeralabs.com"),
    ("8", "What Still Needs to Be Done"),
]

for i, (num, label) in enumerate(items):
    row_t = 1.1 + i * 0.72
    add_rect(slide, 0.4, row_t, 0.45, 0.45, fill=GREEN_DARK)
    add_text(slide, num, 0.4, row_t + 0.02, 0.45, 0.42,
             font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, 1.0, row_t + 0.04, 11.8, 0.45,
             font_size=Pt(16), color=CHARCOAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Spark
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "The Spark — How It All Started")
footer(slide, 3)

# Two-column layout
add_rect(slide, 0.3, 1.1, 6.1, 5.7, fill=WHITE)
add_rect(slide, 6.8, 1.1, 6.1, 5.7, fill=WHITE)

# Left column
txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.7), Inches(5.4))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "The Original Goal"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = GREEN_DARK

add_para(tf, "Build a personal voice assistant — an app I could simply talk to.", Pt(15))
add_para(tf, "JustTalk — a native macOS app powered by Whisper for real-time transcription and voice commands.", Pt(15))
add_para(tf, "It worked beautifully… in English.", Pt(15), italic=True, color=GREEN_MID)

add_para(tf, "The Turning Point", Pt(18), bold=True, color=GREEN_DARK, space_before=Pt(14))
add_para(tf, "\"Can it understand me in Yoruba?\"", Pt(15), italic=True, color=GOLD)
add_para(tf, "\"What about Hausa?\"", Pt(15), italic=True, color=GOLD)
add_para(tf, "Friends and colleagues kept asking the same question — and Whisper had no good answer.", Pt(15))

# Right column
txb2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.2), Inches(5.7), Inches(5.4))
txb2.word_wrap = True
tf2 = txb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.LEFT
run2 = p2.add_run()
run2.text = "What This Revealed"
run2.font.size = Pt(18)
run2.font.bold = True
run2.font.color.rgb = GREEN_DARK

add_para(tf2, "OpenAI Whisper is trained on 680,000 hours of multilingual audio — but Nigerian languages are near-absent.", Pt(15))
add_para(tf2, "Over 220 million Nigerians speak languages that global AI systems effectively cannot hear.", Pt(15))
add_para(tf2, "If AI cannot understand you in your own language, is it really intelligent — or just selectively intelligent?", Pt(15), italic=True, color=GREEN_MID)

add_para(tf2, "→ A personal app project became a research mission.", Pt(16), bold=True, color=GREEN_DARK, space_before=Pt(16))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Problem
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "The Problem — Speech Recognition That Doesn't Know You")
footer(slide, 4)

stats = [
    ("220M+",   "Nigerians whose languages\nare absent from mainstream AI"),
    ("680K hrs", "Whisper training data\n— < 0.1% Nigerian languages"),
    ("4",        "Target languages:\nHausa · Yoruba · Igbo · Pidgin"),
    ("~32K",     "Total samples available\nacross all four languages"),
]

for i, (stat, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.45
    y = 1.15 + row * 2.85
    add_rect(slide, x, y, 6.0, 2.55, fill=WHITE)
    add_rect(slide, x, y, 6.0, 0.55, fill=GREEN_DARK)
    add_text(slide, stat, x + 0.1, y + 0.05, 5.8, 0.5,
             font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.2, y + 0.65, 5.6, 1.8,
             font_size=Pt(16), color=CHARCOAL, align=PP_ALIGN.CENTER)

add_text(slide,
         "The gap is not just technical — it is a question of inclusion. "
         "If intelligent systems cannot understand the world's speakers equally, they serve only a privileged few.",
         0.4, 6.6, 12.5, 0.55,
         font_size=Pt(13), italic=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Technical Approach
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Technical Approach — Fine-Tuning Whisper with LoRA")
footer(slide, 5)

# Pipeline arrow flow
stages = [
    ("Raw Audio\n(16 kHz)", "Input"),
    ("Log-Mel\nSpectrogram", "Feature Extraction"),
    ("Whisper\nEncoder", "Frozen Base"),
    ("LoRA\nAdapters", "Trainable ✓"),
    ("Whisper\nDecoder", "Fine-Tuned"),
    ("Transcribed\nText", "Output"),
]

box_w, box_h = 1.7, 1.3
start_x = 0.4
y_top = 1.25
for i, (main, sub) in enumerate(stages):
    x = start_x + i * 2.1
    color = GREEN_DARK if "LoRA" in main else (GREEN_MID if "Frozen" not in sub else RGBColor(0xAA, 0xAA, 0xAA))
    add_rect(slide, x, y_top, box_w, box_h, fill=color)
    add_text(slide, main, x + 0.05, y_top + 0.1, box_w - 0.1, 0.8,
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x + 0.05, y_top + 0.92, box_w - 0.1, 0.35,
             font_size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_text(slide, "→", x + box_w + 0.05, y_top + 0.35, 0.35, 0.6,
                 font_size=Pt(22), bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

# LoRA key facts
txb = slide.shapes.add_textbox(Inches(0.4), Inches(2.85), Inches(12.5), Inches(3.9))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Why LoRA (Low-Rank Adaptation)?"
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = GREEN_DARK

bullets = [
    "Only 0.273% of model parameters are trained — 196,608 trainable vs 72 million total weights",
    "Full fine-tuning on a 244M-parameter model requires 100s of GB of GPU memory; LoRA runs on a MacBook",
    "Rank-8 matrices injected into Whisper's attention layers (query & value projections)",
    "8-bit quantization (BitsAndBytesConfig) halves memory footprint without accuracy loss",
    "SpecAugment data augmentation (frequency + time masking) reduces overfitting by 20–30% WER",
    "Effective batch size of 16 via gradient accumulation (4 steps × 4 samples)",
]
for b in bullets:
    add_para(tf, f"  •  {b}", Pt(14), color=CHARCOAL, space_before=Pt(5))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Datasets
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Datasets — Sources and Distribution")
footer(slide, 6)

# Table header
headers = ["Language", "Source", "Total Samples", "Train (70%)", "Val (15%)", "Test (15%)"]
col_widths = [1.8, 2.5, 1.7, 1.55, 1.35, 1.35]
col_starts = [0.3]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

row_h = 0.52
header_y = 1.1
for j, (hdr, x, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_rect(slide, x, header_y, w - 0.04, row_h, fill=GREEN_DARK)
    add_text(slide, hdr, x + 0.05, header_y + 0.06, w - 0.1, row_h - 0.1,
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows = [
    ["Hausa",           "Mozilla Common Voice", "27,601", "19,320", "4,140", "4,141"],
    ["Yoruba",          "Google WaxalNLP + Mozilla", "~4,807", "~3,365", "~721", "~721"],
    ["Igbo",            "Mozilla Common Voice", "2,401",  "1,680",  "360",   "361"],
    ["Nigerian Pidgin", "Mozilla Common Voice", "38",     "26",     "6",     "6"],
]
alt = [WHITE, LIGHT_GREEN, WHITE, LIGHT_GREEN]
for ri, (row_data, bg) in enumerate(zip(rows, alt)):
    ry = header_y + row_h + ri * row_h
    for j, (val, x, w) in enumerate(zip(row_data, col_starts, col_widths)):
        add_rect(slide, x, ry, w - 0.04, row_h - 0.02, fill=bg)
        color = GREEN_DARK if j == 0 else CHARCOAL
        add_text(slide, val, x + 0.05, ry + 0.06, w - 0.1, row_h - 0.1,
                 font_size=Pt(13), bold=(j == 0), color=color, align=PP_ALIGN.CENTER)

# Note below table
add_rect(slide, 0.3, 3.45, 12.5, 0.9, fill=RGBColor(0xFF, 0xF9, 0xE5))
add_text(slide,
         "⚠  Dataset Rebalancing: The default Mozilla splits locked ~39% of Hausa data in test.tsv, starving the model.\n"
         "    Solution: All splits merged and redistributed using a strict 70/15/15 algorithmic split, recovering ~2,500 training samples.",
         0.45, 3.5, 12.2, 0.85, font_size=Pt(13), color=CHARCOAL)

add_text(slide,
         "Data is still not enough — this is why we built the crowdsourcing platform (see slide 9)",
         0.3, 4.5, 12.5, 0.45,
         font_size=Pt(13), italic=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

# Source boxes
for src_x, src_label in [(0.3, "Google WaxalNLP"), (4.5, "Mozilla Common Voice"), (8.7, "Local Recordings")]:
    add_rect(slide, src_x, 5.1, 3.8, 0.7, fill=WHITE)
    add_rect(slide, src_x, 5.1, 3.8, 0.2, fill=GREEN_MID)
    add_text(slide, src_label, src_x + 0.1, 5.35, 3.6, 0.4,
             font_size=Pt(14), bold=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

add_text(slide, "Still not enough — more data needed from real Nigerian speakers",
         0.3, 6.0, 12.5, 0.45,
         font_size=Pt(14), italic=True, color=RGBColor(0x99, 0x00, 0x00), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Training Results
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Training Results")
footer(slide, 7)

# WER table
add_text(slide, "Word Error Rate (WER) by Checkpoint — Whisper-Small Full Fine-Tune",
         0.4, 1.1, 12.5, 0.5, font_size=Pt(16), bold=True, color=GREEN_DARK)

wer_headers = ["Checkpoint", "Steps", "WER (%)","Eval Loss", "Notes"]
wer_widths =  [2.2, 1.5, 1.8, 1.8, 5.0]
wer_starts  = [0.4]
for w in wer_widths[:-1]:
    wer_starts.append(wer_starts[-1] + w)

rh = 0.48
hy = 1.65
for j, (h, x, w) in enumerate(zip(wer_headers, wer_starts, wer_widths)):
    add_rect(slide, x, hy, w - 0.04, rh, fill=GREEN_DARK)
    add_text(slide, h, x + 0.05, hy + 0.06, w - 0.1, rh - 0.1,
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

wer_rows = [
    ["Checkpoint-1000", "1,000", "66.34", "1.018", "Early training — model finding its footing"],
    ["Checkpoint-2000", "2,000", "64.42", "0.952", "Steady improvement"],
    ["Checkpoint-3000 ★", "3,000", "64.18", "0.974", "BEST — selected as final model"],
    ["Checkpoint-4000", "4,000", "64.90", "0.996", "Slight regression — overfitting begins"],
]
wer_bgs = [WHITE, LIGHT_GREEN, RGBColor(0xD0, 0xFF, 0xD8), WHITE]
for ri, (row_data, bg) in enumerate(zip(wer_rows, wer_bgs)):
    ry = hy + rh + ri * rh
    for j, (val, x, w) in enumerate(zip(row_data, wer_starts, wer_widths)):
        add_rect(slide, x, ry, w - 0.04, rh - 0.02, fill=bg)
        bold = (ri == 2)
        c = GREEN_DARK if bold else CHARCOAL
        add_text(slide, val, x + 0.05, ry + 0.06, w - 0.1, rh - 0.1,
                 font_size=Pt(12), bold=bold, color=c, align=PP_ALIGN.LEFT if j == 4 else PP_ALIGN.CENTER)

# Insight boxes
insights = [
    ("Best WER Achieved", "64.18%", GREEN_DARK),
    ("Training Loss Start", "4.10  →  0.58", GREEN_MID),
    ("Overfitting Horizon", "After Epoch 4", GOLD),
]
for ii, (label, val, color) in enumerate(insights):
    ix = 0.4 + ii * 4.3
    add_rect(slide, ix, 5.3, 4.0, 1.45, fill=WHITE)
    add_rect(slide, ix, 5.3, 4.0, 0.38, fill=color)
    add_text(slide, label, ix + 0.1, 5.32, 3.8, 0.35,
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, val, ix + 0.1, 5.72, 3.8, 0.9,
             font_size=Pt(22), bold=True, color=color, align=PP_ALIGN.CENTER)

add_text(slide,
         "Note: WER of 64% reflects the difficulty of low-resource tonal language transcription — "
         "not a model failure. Baseline Whisper (untrained) scores near 100% on these languages.",
         0.4, 6.9, 12.5, 0.45, font_size=Pt(11), italic=True, color=CHARCOAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Challenges
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Challenges — What Makes This Hard")
footer(slide, 8)

challenges = [
    ("Data Scarcity",
     "Pidgin: only 38 samples total.\nYoruba & Igbo: ~2,400 each.\nHausa is best-resourced at 27K — still tiny vs English's millions."),
    ("Tonal Languages & Diacritics",
     "Yoruba/Igbo use pitch to change word meaning.\nDiacritics (ọ, ẹ, ṣ, á) dropped by contributors → WER penalises correct output.\nModel generates ọ, ground truth has o → counted as an error."),
    ("Code-Switching",
     "Nigerians naturally mix Yoruba, Hausa, Igbo, Pidgin & English mid-sentence.\nNo training set captures this fluently.\nTokenizer struggles at language boundaries."),
    ("Dialectal Variation",
     "Igbo alone has 20+ regional dialects.\nA model trained on Lagos Yoruba degrades on Ibadan or Ekiti speech.\nNo benchmark exists for dialect-level evaluation."),
    ("Orthographic Inconsistency",
     "Community contributors type without diacritics on standard keyboards.\nSame word transcribed 5 different ways across the dataset.\nNormalisation is non-trivial and can strip meaning."),
    ("Compute Constraints",
     "Full fine-tuning of Whisper-Small requires 40GB+ GPU.\nLoRA + 8-bit quantisation makes it feasible on consumer hardware.\nBut MPS (Apple Silicon) training is 3–5× slower than A100 GPU."),
]

for i, (title, body) in enumerate(challenges):
    col = i % 2
    row = i // 3
    x = 0.3 + col * 6.5
    y = 1.1 + row * 2.0
    add_rect(slide, x, y, 6.2, 1.85, fill=WHITE)
    add_rect(slide, x, y, 0.12, 1.85, fill=GREEN_DARK)
    add_text(slide, title, x + 0.22, y + 0.08, 5.85, 0.38,
             font_size=Pt(14), bold=True, color=GREEN_DARK)
    add_text(slide, body, x + 0.22, y + 0.46, 5.85, 1.3,
             font_size=Pt(12), color=CHARCOAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Crowdsourcing
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "Crowdsourcing Solution — ngs.primeeralabs.com")
footer(slide, 9)

# Left: text explanation
add_rect(slide, 0.3, 1.1, 7.0, 5.7, fill=WHITE)

txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.6), Inches(5.4))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Data Problem Requires a Community Solution"
run.font.size = Pt(17)
run.font.bold = True
run.font.color.rgb = GREEN_DARK

add_para(tf,
    "Existing datasets (Google WaxalNLP, Mozilla Common Voice) are not enough. "
    "The gap can only be closed by Nigerian speakers themselves.",
    Pt(14), space_before=Pt(8))

add_para(tf, "What the platform does:", Pt(15), bold=True, color=GREEN_DARK, space_before=Pt(12))

bullets = [
    "Speakers record short phrases in their native language",
    "Contributions labelled by language and dialect",
    "Audio + transcription submitted directly to training pipeline",
    "Open to all Nigerian language speakers worldwide",
    "Data used to continuously improve the Whisper fine-tune",
]
for b in bullets:
    add_para(tf, f"  ✓  {b}", Pt(14), color=CHARCOAL, space_before=Pt(5))

add_para(tf, "Website: ngs.primeeralabs.com", Pt(15), bold=True, color=GREEN_MID, space_before=Pt(14))

# Right: website highlight card
add_rect(slide, 7.6, 1.1, 5.4, 5.7, fill=GREEN_DARK)
add_text(slide, "ngs.primeeralabs.com", 7.8, 1.3, 5.0, 0.6,
         font_size=Pt(16), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(slide, 8.0, 2.05, 4.6, 0.04, fill=GOLD)

facts = [
    ("Target", "All 4 Nigerian languages\n+ regional dialects"),
    ("Contribution", "Voice recordings\n+ manual transcriptions"),
    ("Goal", "10,000+ samples per\nlanguage by 2027"),
    ("Impact", "Continuously retrain\nWhisper fine-tune"),
]
for fi, (label, val) in enumerate(facts):
    fy = 2.15 + fi * 1.2
    add_rect(slide, 7.8, fy, 4.6, 1.05, fill=RGBColor(0x00, 0x55, 0x30))
    add_text(slide, label, 7.95, fy + 0.05, 4.3, 0.3,
             font_size=Pt(11), bold=True, color=GOLD)
    add_text(slide, val, 7.95, fy + 0.35, 4.3, 0.65,
             font_size=Pt(13), color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — What Still Needs to Be Done
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
section_header(slide, "What Still Needs to Be Done")
footer(slide, 10)

todo = [
    ("More Data",
     "Grow crowdsourced corpus to 10K+ samples per language.\nPartner with Nigerian universities, radio stations, podcast creators."),
    ("Dialect-Level Models",
     "Train separate adapters per dialect (Ibadan Yoruba vs Lagos Yoruba).\nBenchmark on dialect-specific held-out sets."),
    ("Better Evaluation",
     "Move beyond standard WER — develop diacritic-aware WER.\nCharacter Error Rate (CER) is more meaningful for tonal languages."),
    ("Code-Switching Support",
     "Collect mixed-language (Yoruba–English, Hausa–English) recordings.\nTrain a router model to detect language switches mid-utterance."),
    ("Mobile & Offline Deployment",
     "Quantize model to run on Android/iOS without internet.\nNigeria has limited bandwidth — on-device inference is essential."),
    ("JustTalk Integration",
     "Deploy fine-tuned model directly into JustTalk app.\nReal-world user testing with Nigerian language communities."),
]

for i, (title, body) in enumerate(todo):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.1 + row * 2.0
    add_rect(slide, x, y, 6.2, 1.85, fill=WHITE)
    # numbered circle
    add_rect(slide, x + 0.1, y + 0.55, 0.45, 0.45, fill=GREEN_DARK)
    add_text(slide, str(i + 1), x + 0.1, y + 0.56, 0.45, 0.42,
             font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.65, y + 0.08, 5.4, 0.4,
             font_size=Pt(14), bold=True, color=GREEN_DARK)
    add_text(slide, body, x + 0.65, y + 0.5, 5.4, 1.3,
             font_size=Pt(12), color=CHARCOAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Takeaways
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=GREEN_DARK)
add_rect(slide, 0, 0, 0.35, 7.5, fill=GOLD)
footer(slide, 11)

add_text(slide, "Key Takeaways", 0.65, 0.3, 12.0, 0.75,
         font_size=Pt(30), bold=True, color=WHITE)
add_rect(slide, 0.65, 1.1, 12.0, 0.04, fill=GOLD)

takeaways = [
    ("01", "Language Inclusion is an AI Justice Issue",
     "220M+ Nigerians deserve AI that speaks their language. Building for under-represented languages is not niche — it is foundational."),
    ("02", "Parameter-Efficient Fine-Tuning Democratises AI Research",
     "LoRA lets a single researcher adapt a 244M-parameter model on a laptop. You don't need a data centre to do meaningful AI research."),
    ("03", "Data is the Bottleneck — Community is the Solution",
     "No algorithm can compensate for missing data. The crowdsourcing platform at ngs.primeeralabs.com is how we close the gap together."),
    ("04", "The Work is Not Done — But the Path is Clear",
     "More data, dialect-aware evaluation, mobile deployment, and JustTalk integration are the next milestones."),
]

for i, (num, title, body) in enumerate(takeaways):
    y = 1.25 + i * 1.45
    add_rect(slide, 0.55, y, 0.6, 0.6, fill=GOLD)
    add_text(slide, num, 0.55, y + 0.05, 0.6, 0.55,
             font_size=Pt(16), bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.3, y, 11.5, 1.35, fill=RGBColor(0x00, 0x55, 0x30))
    add_text(slide, title, 1.45, y + 0.05, 11.2, 0.42,
             font_size=Pt(15), bold=True, color=GOLD)
    add_text(slide, body, 1.45, y + 0.48, 11.2, 0.8,
             font_size=Pt(13), color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Closing / Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=GREEN_DARK)
add_rect(slide, 0, 0, 0.35, 7.5, fill=GOLD)
footer(slide, 12)

add_text(slide, "Thank You", 0.65, 0.6, 12.0, 1.1,
         font_size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 2.0, 1.75, 9.0, 0.06, fill=GOLD)

add_text(slide,
         '"If it doesn\'t understand me, is it really intelligent?"',
         0.65, 2.0, 12.0, 0.8,
         font_size=Pt(20), italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(slide,
         "Dr. Abimbola Jeremiah\nAssistant Professor, Silesian University of Technology, Poland",
         0.65, 3.0, 12.0, 0.9,
         font_size=Pt(17), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 2.5, 4.1, 8.0, 0.04, fill=RGBColor(0x00, 0x55, 0x30))

contacts = [
    ("GitHub / Code", "github.com/pianistprogrammer/Whisper-Ng"),
    ("Crowdsourcing Platform", "ngs.primeeralabs.com"),
    ("Base Model", "openai/whisper-small  (Hugging Face)"),
    ("Datasets", "google/WaxalNLP  +  Mozilla Common Voice"),
]
for ci, (label, val) in enumerate(contacts):
    cy = 4.25 + ci * 0.65
    add_text(slide, label + ":", 1.5, cy, 3.5, 0.55,
             font_size=Pt(14), bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    add_text(slide, val, 5.2, cy, 7.5, 0.55,
             font_size=Pt(14), color=WHITE)


# ── Save ──────────────────────────────────────────────────────────────────
out_path = "/Users/I558118/Documents/Projects/Whisper-Ng/Whisper_NG_Conference_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
